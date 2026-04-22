#!/usr/bin/env python3
"""Generate Miccy pitch deck — 15 slides, SME/Executive primary audience."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import sys

# ── Brand colours ──────────────────────────────────────────────────────────
BLACK   = RGBColor(0x08, 0x08, 0x10)
DARK    = RGBColor(0x0d, 0x0d, 0x1a)
SURFACE = RGBColor(0x11, 0x11, 0x20)
GOLD    = RGBColor(0xc9, 0xa2, 0x27)
GOLD_DIM= RGBColor(0x6b, 0x56, 0x14)
WHITE   = RGBColor(0xf4, 0xf2, 0xec)
MUTED   = RGBColor(0x99, 0x96, 0x8e)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, colour=BLACK):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def box(slide, l, t, w, h, fill=None, border=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background() if border is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        font="Plus Jakarta Sans", italic=False):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def serif(slide, text, l, t, w, h, size=36, color=WHITE, align=PP_ALIGN.LEFT, bold=False, italic=False):
    return txt(slide, text, l, t, w, h, size=size, bold=bold, color=color,
               align=align, font="DM Serif Display", italic=italic)


def tag(slide, label, l, t):
    """Gold uppercase tag/eyebrow label."""
    b = box(slide, l, t, 2.0, 0.28, fill=RGBColor(0x1a, 0x16, 0x04), border=GOLD_DIM)
    txt(slide, label.upper(), l+0.08, t+0.03, 1.9, 0.25,
        size=7, bold=True, color=GOLD, align=PP_ALIGN.LEFT)


def gold_line(slide, l, t, w):
    """Thin gold rule."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.02))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def slide_number(slide, n):
    txt(slide, str(n), 12.6, 7.1, 0.5, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def footer_bar(slide):
    """Consistent footer across slides 2–15."""
    box(slide, 0, 7.15, 13.33, 0.35, fill=DARK)
    txt(slide, "Miccy by Arus Digital  ·  arusdigital.com  ·  hello@arusdigital.com",
        0.3, 7.18, 9, 0.28, size=8, color=MUTED)
    txt(slide, "CONFIDENTIAL", 10.5, 7.18, 2.5, 0.28, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
box(s, 0, 0, 13.33, 7.5, fill=BLACK)
# Gold accent bar left
box(s, 0, 1.8, 0.08, 4.0, fill=GOLD)
# Gold dot (logo stand-in)
dot = s.shapes.add_shape(9, Inches(0.7), Inches(0.65), Inches(0.18), Inches(0.18))
dot.fill.solid(); dot.fill.fore_color.rgb = GOLD
dot.line.fill.background()
txt(s, "ARUS DIGITAL", 1.0, 0.6, 5, 0.35, size=9, bold=True, color=GOLD)

serif(s, "Your AI\nIntelligence\nOfficer.", 0.7, 1.5, 7.5, 4.0, size=68, color=WHITE)
txt(s, "Daily 7am intelligence briefs for Malaysian CEOs.\nFinancials. Competitors. Anomalies. In your WhatsApp.",
    0.7, 5.3, 8, 0.9, size=15, color=MUTED)

# Right-side accent
box(s, 9.5, 1.5, 3.5, 3.5, fill=SURFACE, border=GOLD_DIM)
serif(s, "Miccy", 9.7, 2.0, 3.1, 1.0, size=28, color=GOLD)
txt(s, "by Arus Digital", 9.7, 2.9, 3.1, 0.4, size=11, color=MUTED)
gold_line(s, 9.7, 3.5, 3.0)
txt(s, "Malaysia's first 100% AI-operated company", 9.7, 3.7, 3.0, 0.8, size=10, color=MUTED)

txt(s, "STRICTLY CONFIDENTIAL  ·  2026", 0.7, 7.1, 6, 0.3, size=8, color=MUTED)
slide_number(s, 1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "The Problem", 0.5, 0.4)
serif(s, "Malaysian CEOs are\nflying blind.", 0.5, 0.8, 8.5, 2.2, size=44, color=WHITE)

# Customer voice quote
box(s, 0.5, 3.2, 9.0, 1.4, fill=SURFACE, border=GOLD_DIM)
serif(s, '"Saya tahu ada masalah — tapi\nsaya tahu last sekali."', 0.8, 3.35, 8.4, 1.1, size=20, italic=True, color=WHITE)
txt(s, "— How Malaysian CEOs describe their situation", 0.8, 4.45, 8, 0.3, size=10, color=MUTED)

# Three pain points
pains = [
    ("10+ systems", "Your data lives in Xero, HubSpot, Shopify, Slack, GA4… and never speaks to each other."),
    ("Reports lag by days", "By the time your team sends a report, the window has already closed."),
    ("Dashboards ≠ decisions", "BI tools show data. Nobody synthesises it. CEOs still interpret everything alone."),
]
for i, (title, body) in enumerate(pains):
    x = 0.5 + i * 4.25
    box(s, x, 5.0, 3.9, 1.9, fill=SURFACE, border=GOLD_DIM)
    txt(s, title, x+0.2, 5.15, 3.5, 0.4, size=11, bold=True, color=GOLD)
    txt(s, body, x+0.2, 5.6, 3.5, 1.2, size=10, color=MUTED)

slide_number(s, 2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "The Solution", 0.5, 0.4)
serif(s, "Miccy — one brief.\nEvery morning at 7am.", 0.5, 0.8, 8, 2.0, size=44, color=WHITE)
txt(s, "Five things that need your attention. Three options each. One recommendation. In your WhatsApp.",
    0.5, 2.9, 8.5, 0.6, size=14, color=MUTED)

# Three key proof points
proofs = [
    ("Private server", "PDPA compliant. Data never leaves your building."),
    ("48h setup", "No data migration. Works with what you already have."),
    ("WhatsApp-native", "Where you already work. No new app to check."),
]
for i, (title, body) in enumerate(proofs):
    x = 0.5 + i * 4.25
    box(s, x, 3.8, 3.9, 1.4, fill=SURFACE, border=GOLD_DIM)
    txt(s, title, x+0.2, 3.95, 3.5, 0.35, size=12, bold=True, color=GOLD)
    txt(s, body, x+0.2, 4.35, 3.5, 0.75, size=10, color=MUTED)

gold_line(s, 0.5, 5.4, 12.3)
txt(s, "Arus Digital is 100% AI-operated internally. Miccy is the same system we use to run our own business. We are our own proof.",
    0.5, 5.55, 12.3, 0.6, size=11, color=MUTED, italic=False)
slide_number(s, 3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "How It Works", 0.5, 0.4)
serif(s, "Three steps.\nZero IT project.", 0.5, 0.8, 8, 1.8, size=44, color=WHITE)

steps = [
    ("01", "Connect", "Point Miccy at your existing systems. Xero, HubSpot, Shopify, Slack, GA4 — no migration required."),
    ("02", "Watch", "Miccy monitors everything 24/7: financials, operations, market signals, competitor moves. Learns your patterns over time."),
    ("03", "Decide", "Every morning at 7am: 5 prioritised alerts. Each with 3 options, a recommendation, and the data behind it. Via WhatsApp."),
]
for i, (num, title, body) in enumerate(steps):
    x = 0.5 + i * 4.25
    box(s, x, 3.0, 3.9, 3.6, fill=SURFACE, border=GOLD_DIM)
    serif(s, num, x+0.2, 3.15, 1.0, 0.7, size=28, color=GOLD)
    txt(s, title, x+0.2, 3.85, 3.5, 0.4, size=13, bold=True, color=WHITE)
    txt(s, body, x+0.2, 4.3, 3.5, 2.2, size=10, color=MUTED)

slide_number(s, 4)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SAMPLE BRIEF
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "Sample Morning Brief", 0.5, 0.4)
serif(s, "This is what lands\nin your WhatsApp at 7am.", 0.5, 0.8, 9, 1.8, size=40, color=WHITE)

# Mock brief card
box(s, 0.5, 2.8, 12.3, 4.0, fill=DARK, border=GOLD_DIM)
txt(s, "📊  Miccy Morning Brief  ·  Thursday, 7:00am", 0.7, 2.95, 12, 0.35, size=10, bold=True, color=GOLD)
gold_line(s, 0.7, 3.35, 12.0)

alerts = [
    ("🔴  CRITICAL", "Revenue down 28% vs 7-day average. 3 invoices unpaid (RM 340,000). Finance team has not flagged this."),
    ("🟡  WATCH",    "Top competitor launched enterprise plan at 40% lower price. 12 new job postings — expansion signal."),
    ("🟢  OPPORTUNITY", "Industry peer acquired analytics firm last week. Market consolidation signal — first-mover window: 60–90 days."),
]
for i, (level, body) in enumerate(alerts):
    y = 3.5 + i * 1.05
    txt(s, level, 0.7, y, 3.0, 0.3, size=9, bold=True, color=WHITE)
    txt(s, body, 0.7, y+0.3, 11.8, 0.6, size=10, color=MUTED)

slide_number(s, 5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — INTEGRATIONS
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "Integrations", 0.5, 0.4)
serif(s, "Works with what\nyou already have.", 0.5, 0.8, 9, 1.8, size=44, color=WHITE)
txt(s, "40+ integrations. 48h setup. No data migration. No IT project.",
    0.5, 2.7, 10, 0.4, size=13, color=MUTED)

integrations = [
    "Xero", "HubSpot", "Shopify", "Slack", "GA4", "Linear",
    "Notion", "Ahrefs", "SimilarWeb", "QuickBooks", "Stripe",
    "Pipedrive", "Salesforce", "Google Sheets", "Telegram", "+ more",
]
cols = 5
for i, name in enumerate(integrations):
    r, c = divmod(i, cols)
    x = 0.5 + c * 2.55
    y = 3.3 + r * 1.1
    box(s, x, y, 2.3, 0.75, fill=SURFACE, border=GOLD_DIM)
    txt(s, name, x+0.1, y+0.18, 2.1, 0.4, size=11, bold=(name == "+ more"), color=WHITE if name != "+ more" else GOLD, align=PP_ALIGN.CENTER)

slide_number(s, 6)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PDPA & DATA SOVEREIGNTY
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "Data Sovereignty", 0.5, 0.4)
serif(s, "Your data never\nleaves your building.", 0.5, 0.8, 9, 1.8, size=44, color=WHITE)
txt(s, "Every Miccy tier ships with private server hosting. PDPA compliant from day one — not retrofitted.",
    0.5, 2.7, 9.5, 0.5, size=13, color=MUTED)

guarantees = [
    ("🔒  Private Server", "Your data is hosted on a dedicated private server — not shared cloud infrastructure. Zero co-tenancy."),
    ("🇲🇾  PDPA Compliant", "Full compliance with Malaysia's Personal Data Protection Act. Audit-ready from day one."),
    ("🚫  Zero External Transfer", "No data sent to OpenAI, Anthropic, or any cloud AI provider. All inference is local."),
    ("🏢  On-Premise Option", "For BFSI, healthcare, and government: we deploy and manage the AI on your own hardware. Data stays on-premises."),
]
for i, (title, body) in enumerate(guarantees):
    r, c = divmod(i, 2)
    x = 0.5 + c * 6.4
    y = 3.5 + r * 1.8
    box(s, x, y, 6.0, 1.55, fill=SURFACE, border=GOLD_DIM)
    txt(s, title, x+0.2, y+0.15, 5.6, 0.4, size=11, bold=True, color=WHITE)
    txt(s, body, x+0.2, y+0.6, 5.6, 0.85, size=10, color=MUTED)

slide_number(s, 7)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PRICING
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "Pricing", 0.5, 0.4)
serif(s, "Three tiers. One\nintelligence officer.", 0.5, 0.8, 9, 1.8, size=44, color=WHITE)

tiers = [
    {
        "name": "Miccy Ops",
        "price": "RM 2,990",
        "impl": "+ RM 5,000 setup",
        "for": "SME owners & ops leaders",
        "features": ["3 custom intelligence workflows", "Ops digest & team coordination", "Automated invoices, tasks, deadlines", "WhatsApp & Telegram delivery", "Private server hosting"],
    },
    {
        "name": "Miccy Executive",
        "price": "RM 5,990",
        "impl": "+ RM 10,000 setup",
        "for": "COOs, founders, division heads",
        "popular": True,
        "features": ["Everything in Ops", "Morning Intelligence Brief — 7am daily", "10 custom intelligence workflows", "Competitor monitoring & market signals", "Financial anomaly detection"],
    },
    {
        "name": "Miccy Command",
        "price": "RM 12,000+",
        "impl": "+ RM 20–35k setup",
        "for": "Group CEOs, GLCs, conglomerates",
        "features": ["Everything in Executive", "Multi-entity / group consolidation", "Custom competitor monitoring scope", "GLC & conglomerate-grade compliance", "SLA-backed support"],
    },
]
for i, tier in enumerate(tiers):
    x = 0.5 + i * 4.25
    border_col = GOLD if tier.get("popular") else GOLD_DIM
    box(s, x, 3.0, 3.9, 3.85, fill=SURFACE, border=border_col)
    if tier.get("popular"):
        txt(s, "MOST POPULAR", x+0.2, 3.1, 3.5, 0.28, size=7, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    serif(s, tier["name"], x+0.2, 3.4, 3.5, 0.6, size=20, color=WHITE)
    serif(s, tier["price"], x+0.2, 4.0, 3.5, 0.5, size=24, color=GOLD)
    txt(s, tier["impl"], x+0.2, 4.5, 3.5, 0.28, size=9, color=MUTED)
    txt(s, "For " + tier["for"], x+0.2, 4.82, 3.5, 0.28, size=9, italic=True, color=MUTED)
    gold_line(s, x+0.2, 5.15, 3.5)
    for j, feat in enumerate(tier["features"]):
        txt(s, "· " + feat, x+0.2, 5.25 + j*0.45, 3.5, 0.38, size=9, color=MUTED)

txt(s, "Setup fee = onboarding investment. ROI in the first anomaly we catch.",
    0.5, 6.85, 12, 0.28, size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
slide_number(s, 8)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — WHO IT'S FOR
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "Who It's For", 0.5, 0.4)
serif(s, "Built for leaders who\nneed to know first.", 0.5, 0.8, 9, 1.8, size=44, color=WHITE)

fits = [
    ("✅  Great fit", [
        "Founder-CEO of an SME with 10–500 staff",
        "COO managing multiple product lines or brands",
        "Growing company with Xero, HubSpot, or Shopify already running",
        "Business that competes on speed of decision",
        "Leadership tired of finding out bad news last",
    ]),
    ("❌  Not a fit", [
        "Pre-revenue startup with no operations to monitor",
        "Business with no digital systems (no CRM, no accounting software)",
        "CEO who doesn't use WhatsApp or Telegram",
        "Organisation where 'reports are fine' and leadership is happy with status quo",
    ]),
]
for i, (label, items) in enumerate(fits):
    x = 0.5 + i * 6.5
    box(s, x, 3.0, 6.0, 3.9, fill=SURFACE, border=GOLD_DIM)
    txt(s, label, x+0.25, 3.15, 5.5, 0.4, size=12, bold=True, color=GOLD)
    for j, item in enumerate(items):
        txt(s, item, x+0.25, 3.65 + j * 0.58, 5.5, 0.5, size=10, color=MUTED)

slide_number(s, 9)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ARUS DIGITAL STORY
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "About Arus Digital", 0.5, 0.4)
serif(s, "We are our own\ncase study.", 0.5, 0.8, 9, 1.8, size=44, color=WHITE)
txt(s, "Malaysia's first 100% AI-operated company. Every function runs through AI agents. No analysts. No manual reports. Just results.",
    0.5, 2.75, 10, 0.6, size=13, color=MUTED)

stack = [
    ("Intelligence", "Miccy", "Daily business brief"),
    ("Orchestration", "Paperclip", "AI org chart"),
    ("Development", "Claude Code", "All engineering"),
    ("Research", "AI Agents", "Competitor & market intel"),
    ("Infrastructure", "OpenClaw", "Servers & deployments"),
    ("Strategy", "You + AI", "Human direction, AI execution"),
]
for i, (func, tool, desc) in enumerate(stack):
    r, c = divmod(i, 3)
    x = 0.5 + c * 4.25
    y = 3.6 + r * 1.65
    box(s, x, y, 3.9, 1.45, fill=SURFACE, border=GOLD_DIM)
    txt(s, func.upper(), x+0.2, y+0.12, 3.5, 0.28, size=7, bold=True, color=GOLD)
    txt(s, tool, x+0.2, y+0.44, 3.5, 0.4, size=13, bold=True, color=WHITE)
    txt(s, desc, x+0.2, y+0.88, 3.5, 0.45, size=10, color=MUTED)

slide_number(s, 10)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ENTERPRISE ADDENDUM (Command tier)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "Enterprise — Miccy Command", 0.5, 0.4)
serif(s, "For group CEOs,\nGLCs, and conglomerates.", 0.5, 0.8, 9, 1.8, size=40, color=WHITE)
txt(s, "Miccy Command delivers multi-entity intelligence — one consolidated brief across all your subsidiaries, brands, and divisions.",
    0.5, 2.7, 10, 0.5, size=13, color=MUTED)

enterprise = [
    ("Multi-Entity Consolidation", "One brief. All subsidiaries. One view of what's happening across your group — without hiring a data team."),
    ("GLC-Grade Compliance", "Designed for government-linked companies and regulated entities. Full PDPA coverage. Audit-ready reporting."),
    ("Custom Competitor Monitoring", "Track any combination of local and regional competitors — pricing moves, hiring signals, product launches, sentiment."),
    ("Dedicated Onboarding Team", "White-glove setup with a dedicated team. Custom workflows built for your group's reporting structure."),
]
for i, (title, body) in enumerate(enterprise):
    r, c = divmod(i, 2)
    x = 0.5 + c * 6.4
    y = 3.5 + r * 1.8
    box(s, x, y, 6.0, 1.55, fill=SURFACE, border=GOLD_DIM)
    txt(s, title, x+0.2, y+0.15, 5.6, 0.4, size=11, bold=True, color=WHITE)
    txt(s, body, x+0.2, y+0.6, 5.6, 0.85, size=10, color=MUTED)

txt(s, "From RM 12,000/month  ·  Implementation: RM 20,000–35,000  ·  SLA-backed support",
    0.5, 6.85, 12, 0.28, size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
slide_number(s, 11)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ON-PREMISE LLM
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "On-Premise AI", 0.5, 0.4)
serif(s, "AI that runs inside\nyour walls.", 0.5, 0.8, 9, 1.8, size=44, color=WHITE)
txt(s, "For enterprises where data cannot leave the building. We deploy and manage open-source AI models on your own infrastructure.",
    0.5, 2.7, 9.5, 0.5, size=13, color=MUTED)

models = [
    ("Qwen3.5", "Best for SMEs", "Multilingual — BM, Chinese, English. Apache 2.0, free commercial use. Best price-to-performance for private LLM.", "RM 15,000–30,000"),
    ("Nemotron-Cascade 2", "Best for enterprise agents", "30B MoE, only 3B active params. Runs on modest hardware. IMO 2025 gold-medal reasoning.", "RM 50,000–120,000"),
    ("GLM-5.1", "Best for agentic workflows", "45.3 on coding benchmarks — 94% of Claude Opus performance. Strong BM, Chinese, English.", "RM 100,000+"),
]
for i, (name, badge, desc, price) in enumerate(models):
    x = 0.5 + i * 4.25
    box(s, x, 3.5, 3.9, 3.25, fill=SURFACE, border=GOLD_DIM)
    txt(s, badge.upper(), x+0.2, 3.65, 3.5, 0.28, size=7, bold=True, color=GOLD)
    serif(s, name, x+0.2, 3.97, 3.5, 0.6, size=18, color=WHITE)
    txt(s, desc, x+0.2, 4.6, 3.5, 1.3, size=10, color=MUTED)
    gold_line(s, x+0.2, 5.95, 3.5)
    txt(s, price, x+0.2, 6.05, 3.5, 0.35, size=11, bold=True, color=GOLD)

slide_number(s, 12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — TRACTION (placeholder-aware)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "Traction", 0.5, 0.4)
serif(s, "Early momentum.", 0.5, 0.8, 9, 1.5, size=44, color=WHITE)

# Honest pilot-stage positioning
box(s, 0.5, 2.5, 12.3, 2.0, fill=SURFACE, border=GOLD_DIM)
txt(s, "Miccy is in active pilot with select Malaysian companies across e-commerce, professional services, and healthcare.", 0.7, 2.7, 11.8, 0.6, size=14, color=WHITE)
txt(s, "Pilot engagements available now — we work closely with early customers to configure custom intelligence workflows and prove ROI before full commitment.", 0.7, 3.3, 11.8, 0.9, size=11, color=MUTED)

metrics = [
    ("48h", "Average setup to first brief"),
    ("5", "Prioritised alerts per morning brief"),
    ("40+", "Integrations supported"),
    ("100%", "Arus Digital runs on Miccy internally"),
]
for i, (num, label) in enumerate(metrics):
    x = 0.5 + i * 3.2
    box(s, x, 4.75, 2.9, 1.5, fill=RGBColor(0x0d, 0x0d, 0x1a), border=GOLD_DIM)
    serif(s, num, x+0.15, 4.9, 2.6, 0.65, size=32, color=GOLD, align=PP_ALIGN.CENTER)
    txt(s, label, x+0.15, 5.55, 2.6, 0.6, size=9, color=MUTED, align=PP_ALIGN.CENTER)

slide_number(s, 13)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — TEAM
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); footer_bar(s)
tag(s, "Team", 0.5, 0.4)
serif(s, "The humans\nbehind the AI.", 0.5, 0.8, 9, 1.8, size=44, color=WHITE)

box(s, 0.5, 3.0, 12.3, 3.75, fill=SURFACE, border=GOLD_DIM)

# Avatar placeholder
dot2 = s.shapes.add_shape(9, Inches(0.9), Inches(3.4), Inches(1.2), Inches(1.2))
dot2.fill.solid(); dot2.fill.fore_color.rgb = DARK
dot2.line.color.rgb = GOLD_DIM; dot2.line.width = Pt(1)
txt(s, "A", 0.9, 3.4, 1.2, 1.2, size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="DM Serif Display")

serif(s, "Amir Ariff", 2.4, 3.2, 6, 0.65, size=28, color=WHITE)
txt(s, "Founder & CEO · Arus Digital", 2.4, 3.85, 6, 0.35, size=11, bold=True, color=GOLD)
txt(s, "Building Malaysia's first AI-operated company from Kuala Lumpur.\nMiccy is the intelligence officer he needed as a CEO — so he built it.",
    2.4, 4.28, 9.8, 0.8, size=12, color=MUTED)

txt(s, "💬 wa.me/60139844412", 2.4, 5.2, 4, 0.35, size=11, color=WHITE)
txt(s, "✉️ hello@arusdigital.com", 6.5, 5.2, 4, 0.35, size=11, color=WHITE)
txt(s, "🌐 arusdigital.com", 2.4, 5.6, 4, 0.35, size=11, color=WHITE)

slide_number(s, 14)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CTA
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
box(s, 0, 0, 13.33, 7.5, fill=BLACK)
box(s, 0, 0.08, 13.33, 0.08, fill=GOLD)  # Gold top bar

serif(s, "Ready to run on AI?", 0.7, 1.2, 11.9, 1.4, size=58, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Book a Miccy pilot — we'll show you what your first brief looks like before you commit.",
    1.0, 2.9, 11.3, 0.55, size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Primary CTA box
box(s, 2.5, 3.7, 8.3, 1.0, fill=GOLD, border=None)
txt(s, "💬  WhatsApp us now  ·  wa.me/60139844412", 2.5, 3.85, 8.3, 0.65,
    size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

txt(s, "Or email: hello@arusdigital.com  ·  arusdigital.com",
    1.0, 5.0, 11.3, 0.45, size=13, color=MUTED, align=PP_ALIGN.CENTER)

gold_line(s, 2.0, 5.7, 9.3)
txt(s, "Miccy by Arus Digital  ·  Malaysia's first 100% AI-operated company  ·  PDPA Compliant  ·  Private server hosting",
    1.0, 5.85, 11.3, 0.45, size=9, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "STRICTLY CONFIDENTIAL  ·  2026", 1.0, 7.1, 11.3, 0.3, size=8, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
out = "/Users/amirariff/projects/arusdigital-web/miccy-pitch-deck.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print(f"  {len(prs.slides)} slides")
