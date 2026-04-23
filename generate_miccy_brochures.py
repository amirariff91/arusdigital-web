"""
Miccy Tier Brochures — 1-pager per tier (Ops, Executive, Command)
Output: miccy-brochure-ops.pptx, miccy-brochure-executive.pptx, miccy-brochure-command.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand ─────────────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x08, 0x08, 0x10)
DARK   = RGBColor(0x0d, 0x0d, 0x1a)
SURF   = RGBColor(0x11, 0x11, 0x20)
GOLD   = RGBColor(0xc9, 0xa2, 0x27)
WHITE  = RGBColor(0xf4, 0xf2, 0xec)
MUTED  = RGBColor(0xa8, 0xa5, 0x98)
TEAL   = RGBColor(0x2d, 0xb8, 0xb8)

# A4 portrait: 8.27 × 11.69 inches
W = Inches(8.27)
H = Inches(11.69)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def add_slide(prs):
    blank = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank)

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background() if line is None else None
    if line is None:
        shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

def txbox(slide, l, t, w, h, text, size, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, wrap=True, font="Plus Jakarta Sans"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font
    return box

def serif_box(slide, l, t, w, h, text, size, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False):
    return txbox(slide, l, t, w, h, text, size, bold=bold, color=color,
                 align=align, italic=italic, font="DM Serif Display")

def divider(slide, t, color=GOLD, opacity=40):
    ln = slide.shapes.add_connector(1, Inches(0.55), t, W - Inches(1.1), t)
    ln.line.color.rgb = color
    ln.line.width = Pt(0.5)

def bullet_item(slide, l, t, w, text, size=9.5, color=MUTED):
    # Arrow + text
    arrow = txbox(slide, l, t, Inches(0.25), Inches(0.3), "→", size,
                  bold=True, color=GOLD)
    body  = txbox(slide, l + Inches(0.25), t, w - Inches(0.25),
                  Inches(0.5), text, size, color=color, wrap=True)
    return body

def badge(slide, l, t, text, bg=GOLD, fg=BLACK):
    r = rect(slide, l, t, Inches(1.6), Inches(0.3), fill=bg)
    txbox(slide, l, t + Inches(0.01), Inches(1.6), Inches(0.3),
          text, 7.5, bold=True, color=fg, align=PP_ALIGN.CENTER)

# ── Shared layout helper ───────────────────────────────────────────────────────
def build_brochure(tier_name, tagline, price, impl_fee, color_accent,
                   features, use_cases, persona_line, cta_line,
                   bm_section=None):
    prs  = new_prs()
    s    = add_slide(prs)

    M  = Inches(0.55)   # margin
    CW = W - 2 * M      # content width

    # ── Full background ──
    rect(s, 0, 0, W, H, fill=BLACK)

    # ── Colour bar top ──
    rect(s, 0, 0, W, Inches(0.06), fill=color_accent)

    # ── Header strip ──────────────────────────────────────────────
    rect(s, 0, Inches(0.06), W, Inches(1.55), fill=DARK)

    # Company name
    txbox(s, M, Inches(0.18), Inches(3), Inches(0.35),
          "ARUS DIGITAL", 7, bold=True, color=color_accent,
          font="Plus Jakarta Sans")

    # Tier badge
    badge(s, W - M - Inches(1.6), Inches(0.18), tier_name,
          bg=color_accent, fg=BLACK)

    # Tier headline
    serif_box(s, M, Inches(0.52), CW, Inches(0.65),
              f"Miccy {tier_name}", 30, color=WHITE)

    # Tagline
    txbox(s, M, Inches(1.1), CW * 0.85, Inches(0.45),
          tagline, 10.5, color=MUTED)

    # ── Divider ──
    divider(s, Inches(1.72))

    # ── Price block ──────────────────────────────────────────────
    px = M
    py = Inches(1.82)

    serif_box(s, px, py, Inches(3.2), Inches(0.55),
              price, 26, color=color_accent)
    txbox(s, px, py + Inches(0.48), Inches(3.2), Inches(0.3),
          "per month  ·  " + impl_fee + " implementation", 8.5, color=MUTED)

    # ── Features ─────────────────────────────────────────────────
    fy = Inches(2.62)
    txbox(s, M, fy, CW, Inches(0.28),
          "WHAT YOU GET", 7.5, bold=True, color=color_accent,
          font="Plus Jakarta Sans")
    fy += Inches(0.32)

    for feat in features:
        bullet_item(s, M, fy, CW, feat, size=9.5, color=MUTED)
        fy += Inches(0.38)

    # ── Use cases ─────────────────────────────────────────────────
    fy += Inches(0.1)
    divider(s, fy)
    fy += Inches(0.18)

    txbox(s, M, fy, CW, Inches(0.28),
          "BUILT FOR", 7.5, bold=True, color=color_accent,
          font="Plus Jakarta Sans")
    fy += Inches(0.3)

    txbox(s, M, fy, CW, Inches(0.45),
          persona_line, 10, color=WHITE, wrap=True)
    fy += Inches(0.48)

    for uc in use_cases:
        bullet_item(s, M, fy, CW, uc, size=9, color=MUTED)
        fy += Inches(0.36)

    # ── BM section (Command only) ─────────────────────────────────
    if bm_section:
        fy += Inches(0.12)
        divider(s, fy, color=GOLD)
        fy += Inches(0.18)
        txbox(s, M, fy, CW, Inches(0.28),
              "DALAM BAHASA MALAYSIA", 7.5, bold=True, color=GOLD,
              font="Plus Jakarta Sans")
        fy += Inches(0.3)
        for line in bm_section:
            bullet_item(s, M, fy, CW, line, size=9, color=MUTED)
            fy += Inches(0.36)

    # ── Objection kills ──────────────────────────────────────────
    fy += Inches(0.1)
    divider(s, fy)
    fy += Inches(0.18)

    txbox(s, M, fy, CW, Inches(0.28),
          "COMMON QUESTIONS", 7.5, bold=True, color=color_accent,
          font="Plus Jakarta Sans")
    fy += Inches(0.3)

    # Inline Q&A mini-cards
    qa_pairs = [
        ("My data is sensitive", "Private server. Zero external AI. PDPA compliant from day one."),
        ("We have dashboards", "Dashboards show. Miccy synthesises. One brief beats ten dashboards."),
        ("48h setup sounds fast", "No data migration. Connects to Xero, HubSpot, Shopify — what you already use."),
    ]
    for q, a in qa_pairs:
        txbox(s, M, fy, CW * 0.38, Inches(0.25),
              q, 8, bold=True, color=WHITE)
        txbox(s, M + CW * 0.38 + Inches(0.1), fy, CW * 0.6, Inches(0.3),
              a, 8, color=MUTED, wrap=True)
        fy += Inches(0.32)

    # ── Bottom CTA strip ─────────────────────────────────────────
    cta_y = H - Inches(1.15)
    rect(s, 0, cta_y, W, Inches(1.15), fill=color_accent)

    # CTA text
    serif_box(s, M, cta_y + Inches(0.14), CW * 0.65, Inches(0.45),
              cta_line, 16, color=BLACK)

    txbox(s, M, cta_y + Inches(0.52), Inches(3.5), Inches(0.28),
          "wa.me/60139844412", 9, bold=True, color=BLACK,
          font="Plus Jakarta Sans")
    txbox(s, M, cta_y + Inches(0.75), Inches(3.5), Inches(0.25),
          "hello@arusdigital.com  ·  arusdigital.com", 8, color=BLACK,
          font="Plus Jakarta Sans")

    # WhatsApp button (right side)
    btn_l = W - M - Inches(2.2)
    btn_t = cta_y + Inches(0.32)
    rect(s, btn_l, btn_t, Inches(2.2), Inches(0.5), fill=BLACK)
    txbox(s, btn_l, btn_t + Inches(0.05), Inches(2.2), Inches(0.4),
          "💬  WhatsApp Us Now →", 9, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER, font="Plus Jakarta Sans")

    return prs


# ══════════════════════════════════════════════════════════════════════════════
# TIER 1 — Miccy Ops
# ══════════════════════════════════════════════════════════════════════════════
ops = build_brochure(
    tier_name   = "Ops",
    tagline     = "Daily AI intelligence brief for SME founders who need to know what's on fire — before their team tells them.",
    price       = "RM 2,990",
    impl_fee    = "RM 5,000",
    color_accent= GOLD,
    features    = [
        "Daily 7am intelligence brief via WhatsApp — 5 prioritised alerts, 3 options each",
        "Revenue anomaly detection: invoice unpaid, dips vs 7-day average, cash flow signals",
        "Operations oversight: task deadlines, team accountability, supply chain alerts",
        "3 custom monitoring workflows tailored to your business",
        "40+ integrations: Xero, HubSpot, Shopify, Slack, GA4 and more",
        "Private server deployment — your data never leaves your infrastructure",
        "48h setup. No data migration. Works with what you already have.",
    ],
    use_cases   = [
        "SME founders managing 10–100 staff across multiple systems",
        "Owner-operators who are last to know when something breaks",
        "Growing companies replacing manual Excel reports and missed WhatsApp alerts",
    ],
    persona_line = "Founder-CEOs and owner-operators of Malaysian SMEs who need synthesis, not more dashboards.",
    cta_line    = "Book a Miccy demo.",
)
ops.save("miccy-brochure-ops.pptx")
print("✓ miccy-brochure-ops.pptx")


# ══════════════════════════════════════════════════════════════════════════════
# TIER 2 — Miccy Executive
# ══════════════════════════════════════════════════════════════════════════════
exec_brochure = build_brochure(
    tier_name   = "Executive",
    tagline     = "For CEOs who need competitor intelligence, market signals, and morning clarity — not just internal metrics.",
    price       = "RM 5,990",
    impl_fee    = "RM 10,000",
    color_accent= TEAL,
    features    = [
        "Everything in Miccy Ops, plus:",
        "Competitor monitoring: pricing changes, hiring signals, product launches, PR moves",
        "Market intelligence: M&A signals, industry shifts, first-mover window detection (60–90 day alerts)",
        "Board-ready summaries: weekly digest formatted for exec review",
        "5 custom monitoring workflows — internal ops + external market watch",
        "Bahasa Malaysia support for local market signals and local competitor context",
        "Telegram delivery option in addition to WhatsApp",
        "Priority support and onboarding with Arus Digital team",
    ],
    use_cases   = [
        "CEOs of RM10M–RM200M revenue companies competing in fast-moving markets",
        "COOs and division heads who need operational + market visibility in one brief",
        "Executives who've missed a competitor move and won't let it happen again",
    ],
    persona_line = "CEOs and COOs who need both internal anomalies AND external market intelligence — synthesised, not scattered.",
    cta_line    = "Lead your market. Start here.",
)
exec_brochure.save("miccy-brochure-executive.pptx")
print("✓ miccy-brochure-executive.pptx")


# ══════════════════════════════════════════════════════════════════════════════
# TIER 3 — Miccy Command (BM + English)
# ══════════════════════════════════════════════════════════════════════════════
cmd_gold = RGBColor(0xe8, 0xc5, 0x4a)  # slightly brighter gold for Command

command = build_brochure(
    tier_name   = "Command",
    tagline     = "Enterprise-grade AI intelligence for group CEOs, GLCs and conglomerates overseeing multiple entities.",
    price       = "RM 12,000+",
    impl_fee    = "RM 20,000–35,000",
    color_accent= cmd_gold,
    features    = [
        "Everything in Miccy Executive, plus:",
        "Multi-entity consolidation: single brief across all subsidiaries and divisions",
        "GLC-grade compliance: PDPA, data residency, audit trail, role-based access control",
        "On-premise LLM option: AI model runs on your servers — zero external data transfer",
        "Custom KPI frameworks mapped to group-level and entity-level targets",
        "Board reporting templates: auto-drafted summaries for governance committees",
        "Unlimited monitoring workflows — configured to your group's operating model",
        "Dedicated implementation team + quarterly business reviews with Arus Digital",
    ],
    use_cases   = [
        "Group CEOs overseeing 3–20 subsidiaries with fragmented data across entities",
        "GLCs and government-linked funds requiring PDPA-strict AI with on-premise deployment",
        "Conglomerates where board-level visibility currently requires a team of analysts",
    ],
    persona_line = "Group CEOs, Group CFOs, and GLC leaders who need consolidated intelligence across entities — securely.",
    cta_line    = "Enterprise AI. Malaysian-built.",
    bm_section  = [
        "Miccy Command direka khas untuk syarikat korporat besar, GLC dan konglomerat Malaysia.",
        "Satu laporan kecerdasan setiap pagi — merangkumi semua anak syarikat, dalam satu mesej WhatsApp.",
        "Pematuhan PDPA sepenuhnya. Data tidak keluar dari infrastruktur anda. Pilihan LLM swasta tersedia.",
        "Pasukan kami bekerja secara langsung dengan pasukan IT dan lembaga pengarah anda.",
    ],
)
command.save("miccy-brochure-command.pptx")
print("✓ miccy-brochure-command.pptx")

print("\nAll 3 brochures generated.")
